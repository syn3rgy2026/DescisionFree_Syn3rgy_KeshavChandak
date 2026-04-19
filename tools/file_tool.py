"""
file_tool.py
------------
Universal file operations toolkit capable of creating, reading, modifying,
and managing files across multiple formats including:
* Text files (.txt, .md)
* Data files (.csv, .json, .xml)
* Excel files (.xlsx)
* PowerPoint presentations (.pptx)
* Word documents (.docx)
* Code files (.py, .js, .cpp, etc.)

Writes: relative paths go under `output/` (config.OUTPUT_FOLDER). Absolute paths are allowed
under the project working directory, that output folder, or the user's home (so Desktop
and similar paths from the user work). Other locations are rejected.
"""

import os
import csv
import json
import logging
from pathlib import Path
from typing import Optional, List, Dict, Any
from smolagents import tool
import config
from tools.scraper_tool import scrape_website, scrape_website_structured
from tools.web_screenshot import capture_website_screenshot
from tools.web_search_tool import search_web, search_news

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger("file_tool")


def _output_directory_abs() -> str:
    """Absolute path to the default artefact directory (project output/)."""
    of = config.OUTPUT_FOLDER.replace("\\", "/").strip()
    if of.startswith("./"):
        of = of[2:].strip("/")
    return str(Path.cwd() / of)


def _write_allowed_roots() -> tuple[str, ...]:
    """Writes are allowed under: default output dir, cwd, user home (Desktop, Documents, …)."""
    out = _output_directory_abs()
    cwd = str(Path.cwd().resolve())
    home = str(Path.home().resolve())
    return tuple(dict.fromkeys((out, cwd, home)))


def _is_under_root(root: str, path: str) -> bool:
    try:
        root_r = Path(root).resolve()
        path_r = Path(path).resolve()
    except OSError:
        return False
    try:
        path_r.relative_to(root_r)
        return True
    except ValueError:
        return path_r == root_r


def _path_allowed_for_write(abs_path: str) -> bool:
    ap = str(Path(abs_path))
    for r in _write_allowed_roots():
        if _is_under_root(r, ap):
            return True
    return False


def _resolve_write_path(filename: str) -> str:
    """
    Resolve a target path for create/write tools.

    - Absolute paths (after ~ / env expand) must stay under project cwd, OUTPUT_FOLDER,
      or the user's home directory so Desktop/Documents requests work.
    - Relative paths are rooted under OUTPUT_FOLDER (not bare cwd), so `report.csv`
      becomes `<project>/output/report.csv`.
    """
    raw = (filename or "").strip()
    if not raw:
        raise ValueError("filename is empty")

    expanded = os.path.expandvars(os.path.expanduser(raw))
    out_abs = _output_directory_abs()

    if Path(expanded).is_absolute():
        full = str(Path(expanded).resolve())
        parent = str(Path(full).parent)
        if parent and parent != full:
            os.makedirs(parent, exist_ok=True)
        if not _path_allowed_for_write(full):
            raise ValueError(
                "Refusing to write outside allowed locations (project output/, cwd, or user home). "
                f"Got: {full}"
            )
        return full

    rel = expanded.replace("\\", "/").lstrip("./")
    out_seg = Path(out_abs).name
    parts = rel.split("/")
    if parts and parts[0] == out_seg:
        rel = "/".join(parts[1:]) if len(parts) > 1 else ""
    target = str(Path(out_abs) / rel) if rel else out_abs
    full = str(Path(target).resolve())
    if not _is_under_root(out_abs, full):
        raise ValueError(f"Invalid relative path (escapes output folder): {filename!r}")
    parent = str(Path(full).parent)
    if parent:
        os.makedirs(parent, exist_ok=True)
    return full


def _resolve_delete_path(path: str) -> str:
    p = os.path.expandvars(os.path.expanduser(path.strip()))
    if not Path(p).is_absolute():
        p = _resolve_write_path(path)
    full = str(Path(p).resolve())
    if not _path_allowed_for_write(full):
        raise ValueError(f"Refusing to delete outside allowed locations: {full}")
    return full

def _log_action(action: str, details: str):
    logger.info(f"{action} | {details}")


def _resolve_write_path_result(filename: str) -> tuple[str | None, str | None]:
    try:
        return _resolve_write_path(filename), None
    except ValueError as e:
        return None, str(e)


@tool
def read_file(path: str) -> str:
    """
    Read and return the contents of a file (Text, Code, etc).

    Args:
        path: Absolute or relative path to the file.

    Returns:
        Full text content of the file.
    """
    rp = os.path.expandvars(os.path.expanduser(path))
    with open(rp, "r", encoding="utf-8") as f:
        return f.read()

@tool
def write_file(filename: str, content: str) -> str:
    """
    Write text or code to a file. If the user asked for Desktop/Documents/home or gave an
    absolute path, use that full path (e.g. ~/Desktop/notes.md). If they did not specify a
    location, a relative name like `report.md` is stored under the project `output/` folder.

    Args:
        filename: Relative name under output/, or path under home/cwd as allowed.
        content: Text content to write.

    Returns:
        Absolute path to the written file.
    """
    full_path, err = _resolve_write_path_result(filename)
    if err:
        return f"ERROR: {err}"
    with open(full_path, "w", encoding="utf-8") as f:
        f.write(content)
    _log_action("Write file", f"Written to {full_path}")
    return full_path

@tool
def append_file(filename: str, content: str) -> str:
    """
    Append content to an existing file (same path rules as write_file).

    Args:
        filename: Target path (relative → under output/).
        content: Text to append.

    Returns:
        Absolute path to the modified file.
    """
    full_path, err = _resolve_write_path_result(filename)
    if err:
        return f"ERROR: {err}"
    with open(full_path, "a", encoding="utf-8") as f:
        f.write(content)
    _log_action("Append file", f"Appended to {full_path}")
    return full_path

@tool
def list_files(directory: str) -> list:
    """
    List all files in the specified directory.

    Args:
        directory: Path to the directory to list.

    Returns:
        Filenames found in the directory.
    """
    return os.listdir(directory)

@tool
def delete_file(path: str) -> bool:
    """
    Delete a file at the given path (same allowed roots as writes).

    Args:
        path: Path to the file to delete.

    Returns:
        True if deletion succeeded.
    """
    try:
        full = _resolve_delete_path(path)
        if os.path.exists(full):
            os.remove(full)
            _log_action("Delete file", f"Deleted {full}")
        return True
    except (ValueError, OSError, PermissionError):
        return False

@tool
def read_json_file(path: str) -> dict:
    """
    Reads a JSON file safely, maintaining valid schema structure.

    Args:
        path: Path to the target JSON file.

    Returns:
        The dictionary representation of the JSON data.
    """
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

@tool
def write_json_file(filename: str, data: dict) -> str:
    """
    Writes valid JSON content to a file inside the output/ directory.
    Maintains clean key-value structure formatting.

    Args:
        filename: Name of the file to create or overwrite (e.g., 'data.json').
        data: Dictionary object to save as JSON.

    Returns:
        Absolute path to the written file.
    """
    full_path, err = _resolve_write_path_result(filename)
    if err:
        return f"ERROR: {err}"
    with open(full_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4)
    _log_action("Write JSON", f"Written to {full_path}")
    return full_path

@tool
def write_csv_file(filename: str, headers: list, rows: list) -> str:
    """
    Creates a CSV file inside the output/ directory ensuring consistent column counts.

    Args:
        filename: Name of the CSV file (e.g., 'data.csv').
        headers: List of column names.
        rows: List of rows, where each row is a list of values.

    Returns:
        Absolute path to the written file.
    """
    full_path, err = _resolve_write_path_result(filename)
    if err:
        return f"ERROR: {err}"
    with open(full_path, mode='w', newline='', encoding='utf-8') as file:
        writer = csv.writer(file)
        writer.writerow(headers)
        writer.writerows(rows)
    _log_action("Write CSV", f"Written to {full_path}")
    return full_path

@tool
def create_excel_file(filename: str, headers: list, rows: list) -> str:
    """
    Creates an Excel (.xlsx) file in the output/ directory maintaining tabular structure.
    Always creates at least one sheet and avoids empty rows.

    Args:
        filename: Name of the standard Excel file (e.g., 'report.xlsx').
        headers: Row headers representing the schema.
        rows: Grid data matching header columns.

    Returns:
        Absolute path to the written file.
    """
    full_path, err = _resolve_write_path_result(filename)
    if err:
        return f"ERROR: {err}"
    try:
        import xlsxwriter
        workbook = xlsxwriter.Workbook(full_path)
        worksheet = workbook.add_worksheet()
        for col_num, header in enumerate(headers):
            worksheet.write(0, col_num, header)
        for row_num, row_data in enumerate(rows):
            for col_num, cell_data in enumerate(row_data):
                worksheet.write(row_num + 1, col_num, cell_data)
        workbook.close()
        _log_action("Create Excel", f"Written to {full_path}")
    except ImportError:
        raise ImportError("xlsxwriter is required for Excel creation. Install with: pip install XlsxWriter")
    return full_path

@tool
def create_word_document(filename: str, title: str, paragraphs: list) -> str:
    """
    Generates a structured Word (.docx) document inside the output/ directory.

    Args:
        filename: Name of the Word document (e.g., 'document.docx').
        title: Bolded heading 1 title.
        paragraphs: List of string paragraphs to insert sequentially.

    Returns:
        Absolute path to the written file.
    """
    full_path, err = _resolve_write_path_result(filename)
    if err:
        return f"ERROR: {err}"
    try:
        import docx
        doc = docx.Document()
        if title:
            doc.add_heading(title, 0)
        for p in paragraphs:
            doc.add_paragraph(p)
        doc.save(full_path)
        _log_action("Create Word", f"Written to {full_path}")
    except ImportError:
        raise ImportError("python-docx is required. Install with: pip install python-docx")
    return full_path

@tool
def create_ppt_presentation(filename: str, slide_title: str, bullet_points: list) -> str:
    """
    Generates a structured PowerPoint (.pptx) presentation inside the output/ directory.
    Creates a title slide and a content slide (ensuring clear bullet organization).

    Args:
        filename: Name of the presentation file (e.g., 'deck.pptx').
        slide_title: Header title for the presentation layout.
        bullet_points: Maximum list of 5 concise text bullet points.

    Returns:
        Absolute path to the written file.
    """
    full_path, err = _resolve_write_path_result(filename)
    if err:
        return f"ERROR: {err}"
    try:
        from pptx import Presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_slide_layout)
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        title.text = slide_title
        subtitle.text = "Generated by Synergy Agent"

        # Content slide with bullets
        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide2.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_title
        tf = body_shape.text_frame
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point

        prs.save(full_path)
        _log_action("Create PPT", f"Written to {full_path}")
    except ImportError:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")
    return full_path

ALL_TOOLS = [
    read_file,
    write_file,
    append_file,
    list_files,
    delete_file,
    read_json_file,
    write_json_file,
    write_csv_file,
    create_excel_file,
    create_word_document,
    create_ppt_presentation,
    scrape_website,
    scrape_website_structured,
    capture_website_screenshot,
    search_web,
    search_news
]
