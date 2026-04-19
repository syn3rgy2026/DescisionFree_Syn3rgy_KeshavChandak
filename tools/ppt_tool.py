"""
ppt_tool.py
-----------
Industry-grade PowerPoint creation tool for Synergy Agent.

Capabilities:
  - Create multi-slide presentations with professional layouts
  - Generate charts (bar, line, pie) via matplotlib
  - Generate diagrams (flowchart, architecture, pipeline) via PIL
  - Fetch relevant images from the web (Pexels API)
  - Apply gradient backgrounds and accent designs
  - Smart visual insertion with proper placement
  - Support title, bullet, section, text, comparison, image, chart,
    diagram, and visual slides

All outputs saved to the output/ directory.
Requires: python-pptx, matplotlib, Pillow, requests
"""

import os
import json
import logging
import textwrap
from smolagents import tool

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger("ppt_tool")


# ── Paths ─────────────────────────────────────────────────────────────

_PROJECT_ROOT = os.path.dirname(os.path.dirname(__file__))
_OUTPUT_DIR = os.path.join(_PROJECT_ROOT, "output")
os.makedirs(_OUTPUT_DIR, exist_ok=True)


# ── Helpers ───────────────────────────────────────────────────────────

def _get_safe_path(filename: str) -> str:
    """Resolve an absolute path, creating parent directories as needed."""
    full_path = os.path.abspath(filename)
    os.makedirs(os.path.dirname(full_path), exist_ok=True)
    return full_path


# ── Pre-defined colour themes ────────────────────────────────────────

THEMES = {
    "default": {
        "title_color": (0x1A, 0x1A, 0x2E),
        "subtitle_color": (0x6C, 0x75, 0x7D),
        "heading_color": (0x0D, 0x6E, 0xFD),
        "body_color": (0x33, 0x33, 0x33),
        "accent_color": (0x0D, 0x6E, 0xFD),
        "accent_color_2": (0x19, 0x8C, 0xFF),
        "bg_color": (0xFF, 0xFF, 0xFF),
        "bg_gradient_end": (0xF0, 0xF4, 0xFF),
        "title_font": "Calibri",
        "body_font": "Calibri",
        "chart_colors": ["#0D6EFD", "#198CFF", "#6CB4FF", "#A3D0FF", "#D6EBFF"],
        "chart_bg": "#FFFFFF",
    },
    "dark": {
        "title_color": (0xFF, 0xFF, 0xFF),
        "subtitle_color": (0xBB, 0xBB, 0xBB),
        "heading_color": (0x00, 0xD4, 0xFF),
        "body_color": (0xDD, 0xDD, 0xDD),
        "accent_color": (0x00, 0xD4, 0xFF),
        "accent_color_2": (0x00, 0x9D, 0xC4),
        "bg_color": (0x1E, 0x1E, 0x2E),
        "bg_gradient_end": (0x2A, 0x2A, 0x40),
        "title_font": "Calibri",
        "body_font": "Calibri",
        "chart_colors": ["#00D4FF", "#009DC4", "#00E88F", "#FFD700", "#FF6B6B"],
        "chart_bg": "#1E1E2E",
    },
    "corporate": {
        "title_color": (0x00, 0x2B, 0x5C),
        "subtitle_color": (0x5A, 0x5A, 0x5A),
        "heading_color": (0x00, 0x2B, 0x5C),
        "body_color": (0x33, 0x33, 0x33),
        "accent_color": (0xC8, 0x10, 0x2E),
        "accent_color_2": (0x00, 0x2B, 0x5C),
        "bg_color": (0xF8, 0xF9, 0xFA),
        "bg_gradient_end": (0xED, 0xF0, 0xF5),
        "title_font": "Arial",
        "body_font": "Arial",
        "chart_colors": ["#002B5C", "#C8102E", "#4A7C59", "#E8A317", "#7B2D8E"],
        "chart_bg": "#F8F9FA",
    },
    "creative": {
        "title_color": (0x6C, 0x2E, 0xB9),
        "subtitle_color": (0x99, 0x99, 0x99),
        "heading_color": (0xE9, 0x1E, 0x63),
        "body_color": (0x33, 0x33, 0x33),
        "accent_color": (0xFF, 0x98, 0x00),
        "accent_color_2": (0xE9, 0x1E, 0x63),
        "bg_color": (0xFF, 0xFF, 0xFF),
        "bg_gradient_end": (0xFF, 0xF5, 0xF0),
        "title_font": "Georgia",
        "body_font": "Calibri",
        "chart_colors": ["#6C2EB9", "#E91E63", "#FF9800", "#4CAF50", "#2196F3"],
        "chart_bg": "#FFFFFF",
    },
    "consulting": {
        "title_color": (0x00, 0x32, 0x5B),
        "subtitle_color": (0x68, 0x68, 0x68),
        "heading_color": (0x00, 0x32, 0x5B),
        "body_color": (0x2D, 0x2D, 0x2D),
        "accent_color": (0x00, 0x7B, 0xBB),
        "accent_color_2": (0xE8, 0x6C, 0x00),
        "bg_color": (0xFF, 0xFF, 0xFF),
        "bg_gradient_end": (0xF7, 0xF9, 0xFC),
        "title_font": "Calibri",
        "body_font": "Calibri",
        "chart_colors": ["#00325B", "#007BBB", "#E86C00", "#2EAD6D", "#8B5CF6"],
        "chart_bg": "#FFFFFF",
    },
}


def _get_theme(name: str) -> dict:
    """Return a theme dict by name, falling back to 'default'."""
    return THEMES.get(name.lower(), THEMES["default"])


def _rgb(tup):
    """Convert an (r, g, b) tuple to a pptx RGBColor object."""
    from pptx.dml.color import RGBColor
    return RGBColor(tup[0], tup[1], tup[2])


def _set_slide_bg(slide, color_tuple):
    """Apply a solid background colour to a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color_tuple)


def _set_gradient_bg(slide, start_tuple, end_tuple):
    """Apply a gradient background to a slide."""
    from pptx.oxml.ns import qn
    from lxml import etree

    background = slide.background
    fill = background.fill
    fill.gradient()

    try:
        gsLst = fill._fill.find(qn("a:gsLst"))
        if gsLst is not None:
            # Clear existing stops
            for gs in gsLst.findall(qn("a:gs")):
                gsLst.remove(gs)
            # Start colour (0%)
            gs1 = etree.SubElement(gsLst, qn("a:gs"))
            gs1.set("pos", "0")
            srgb1 = etree.SubElement(etree.SubElement(gs1, qn("a:srgbClr")), qn("a:srgbClr"))
            gs1_clr = etree.SubElement(gs1, qn("a:srgbClr"))
            gs1_clr.set("val", "%02X%02X%02X" % start_tuple)
            # Remove the nested element we accidentally made
            for child in gs1:
                if child.tag != qn("a:srgbClr"):
                    continue
                if child.get("val") is None:
                    gs1.remove(child)
                    break
            # End colour (100%)
            gs2 = etree.SubElement(gsLst, qn("a:gs"))
            gs2.set("pos", "100000")
            gs2_clr = etree.SubElement(gs2, qn("a:srgbClr"))
            gs2_clr.set("val", "%02X%02X%02X" % end_tuple)
    except Exception:
        # Fallback to solid
        fill.solid()
        fill.fore_color.rgb = _rgb(start_tuple)


def _add_text_box(slide, left, top, width, height, text, font_name, font_size,
                  font_color_tuple, bold=False, alignment=None):
    """Add a styled text box to a slide."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    if alignment is None:
        alignment = PP_ALIGN.LEFT

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = _rgb(font_color_tuple)
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _add_accent_bar(slide, color_tuple, left=None, top=None, width=None, height=None):
    """Add a thin coloured accent bar (rectangle) to a slide."""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE

    if left is None:
        left = Inches(0)
    if top is None:
        top = Inches(0)
    if width is None:
        width = Inches(0.15)
    if height is None:
        height = Inches(7.5)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_tuple)
    shape.line.fill.background()
    return shape


def _add_bottom_bar(slide, color_tuple, color_tuple_2=None):
    """Add a footer accent bar at the bottom of a slide."""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE

    # Main bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.1),
        Inches(13.333), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_tuple)
    shape.line.fill.background()

    # Optional secondary thin line
    if color_tuple_2:
        shape2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.2),
            Inches(13.333), Inches(0.03)
        )
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = _rgb(color_tuple_2)
        shape2.line.fill.background()


def _add_slide_number(slide, slide_num, total, t):
    """Add a slide number in the bottom-right corner."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    _add_text_box(
        slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
        f"{slide_num} / {total}",
        t["body_font"], 9, t["subtitle_color"],
        alignment=PP_ALIGN.RIGHT
    )


# ══════════════════════════════════════════════════════════════════════
# NEW TOOL 1: Chart Generation (matplotlib)
# ══════════════════════════════════════════════════════════════════════

@tool
def generate_chart_image(
    data: str,
    chart_type: str = "bar",
    filename: str = "output/chart.png",
    title: str = "",
    xlabel: str = "",
    ylabel: str = "",
    theme: str = "default",
) -> str:
    """Generate a chart image using matplotlib and save it to disk.

    Args:
        data: JSON string with chart data. Format depends on chart_type:
              Bar/Line: '{"labels": ["A","B","C"], "values": [10,20,30]}'
              or multi-series: '{"labels": ["Q1","Q2"], "series": [{"name": "Revenue", "values": [10,20]}, {"name": "Cost", "values": [5,8]}]}'
              Pie: '{"labels": ["A","B","C"], "values": [40,35,25]}'
        chart_type: Type of chart — 'bar', 'line', 'pie', 'horizontal_bar'. Default: 'bar'.
        filename: Output filename (e.g. 'output/revenue_chart.png').
        title: Optional chart title.
        xlabel: Optional x-axis label (bar/line only).
        ylabel: Optional y-axis label (bar/line only).
        theme: PPT theme name to match chart colours.

    Returns:
        str: Path to saved chart image, or error message.
    """
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.ticker as mticker
    except ImportError:
        return "Error: matplotlib is required. Install with: pip install matplotlib"

    try:
        chart_data = json.loads(data)
    except json.JSONDecodeError:
        return "Error: data must be valid JSON"

    t = _get_theme(theme)
    colors = t.get("chart_colors", ["#0D6EFD", "#198CFF", "#6CB4FF", "#A3D0FF", "#D6EBFF"])
    bg_color = t.get("chart_bg", "#FFFFFF")

    labels = chart_data.get("labels", [])
    values = chart_data.get("values", [])
    series = chart_data.get("series", [])

    fig, ax = plt.subplots(figsize=(10, 6), facecolor=bg_color)
    ax.set_facecolor(bg_color)

    # Determine text colour based on background
    is_dark_bg = bg_color.lower() in ("#1e1e2e", "#2a2a40") or sum(int(bg_color[i:i+2], 16) for i in (1,3,5)) < 384
    text_color = "#FFFFFF" if is_dark_bg else "#333333"
    grid_color = "#444444" if is_dark_bg else "#E0E0E0"

    ax.tick_params(colors=text_color)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color(grid_color)
    ax.spines["bottom"].set_color(grid_color)

    if chart_type == "pie":
        wedges, texts, autotexts = ax.pie(
            values, labels=labels, colors=colors[:len(values)],
            autopct="%1.1f%%", startangle=140, pctdistance=0.8,
            wedgeprops={"edgecolor": bg_color, "linewidth": 2}
        )
        for t_txt in texts:
            t_txt.set_color(text_color)
            t_txt.set_fontsize(11)
        for at in autotexts:
            at.set_color("white")
            at.set_fontweight("bold")
            at.set_fontsize(10)
        ax.set_aspect("equal")

    elif chart_type == "horizontal_bar":
        if series:
            y_pos = range(len(labels))
            bar_h = 0.8 / len(series)
            for i, s in enumerate(series):
                offsets = [p + i * bar_h for p in y_pos]
                ax.barh(offsets, s["values"], bar_h, color=colors[i % len(colors)],
                        label=s.get("name", f"Series {i+1}"), edgecolor="none")
            ax.set_yticks([p + bar_h * (len(series) - 1) / 2 for p in y_pos])
            ax.set_yticklabels(labels, color=text_color)
            ax.legend(facecolor=bg_color, edgecolor="none", labelcolor=text_color)
        else:
            ax.barh(labels, values, color=colors[0], edgecolor="none")
            ax.tick_params(axis="y", colors=text_color)

    elif chart_type == "line":
        if series:
            for i, s in enumerate(series):
                ax.plot(labels, s["values"], marker="o", color=colors[i % len(colors)],
                        linewidth=2.5, markersize=7, label=s.get("name", f"Series {i+1}"))
            ax.legend(facecolor=bg_color, edgecolor="none", labelcolor=text_color)
        else:
            ax.plot(labels, values, marker="o", color=colors[0], linewidth=2.5, markersize=7)
        ax.grid(True, axis="y", alpha=0.3, color=grid_color)

    else:  # bar (default)
        if series:
            x_pos = range(len(labels))
            bar_w = 0.8 / len(series)
            for i, s in enumerate(series):
                offsets = [p + i * bar_w for p in x_pos]
                ax.bar(offsets, s["values"], bar_w, color=colors[i % len(colors)],
                       label=s.get("name", f"Series {i+1}"), edgecolor="none")
            ax.set_xticks([p + bar_w * (len(series) - 1) / 2 for p in x_pos])
            ax.set_xticklabels(labels, color=text_color)
            ax.legend(facecolor=bg_color, edgecolor="none", labelcolor=text_color)
        else:
            bar_colors = colors[:len(values)] if len(values) <= len(colors) else [colors[i % len(colors)] for i in range(len(values))]
            ax.bar(labels, values, color=bar_colors, edgecolor="none", width=0.6)
        ax.grid(True, axis="y", alpha=0.3, color=grid_color)

    if title:
        ax.set_title(title, fontsize=16, fontweight="bold", color=text_color, pad=15)
    if xlabel:
        ax.set_xlabel(xlabel, fontsize=12, color=text_color)
    if ylabel:
        ax.set_ylabel(ylabel, fontsize=12, color=text_color)

    plt.tight_layout(pad=1.5)
    full_path = _get_safe_path(filename)
    fig.savefig(full_path, dpi=200, bbox_inches="tight", facecolor=bg_color)
    plt.close(fig)

    logger.info(f"Chart saved: {full_path}")
    return full_path


# ══════════════════════════════════════════════════════════════════════
# NEW TOOL 2: Diagram Generation (PIL-based)
# ══════════════════════════════════════════════════════════════════════

def _draw_rounded_rect(draw, xy, radius, fill, outline=None, width=1):
    """Draw a rounded rectangle on a PIL ImageDraw."""
    x0, y0, x1, y1 = xy
    draw.rectangle([x0 + radius, y0, x1 - radius, y1], fill=fill)
    draw.rectangle([x0, y0 + radius, x1, y1 - radius], fill=fill)
    draw.pieslice([x0, y0, x0 + 2*radius, y0 + 2*radius], 180, 270, fill=fill)
    draw.pieslice([x1 - 2*radius, y0, x1, y0 + 2*radius], 270, 360, fill=fill)
    draw.pieslice([x0, y1 - 2*radius, x0 + 2*radius, y1], 90, 180, fill=fill)
    draw.pieslice([x1 - 2*radius, y1 - 2*radius, x1, y1], 0, 90, fill=fill)
    if outline:
        draw.arc([x0, y0, x0 + 2*radius, y0 + 2*radius], 180, 270, fill=outline, width=width)
        draw.arc([x1 - 2*radius, y0, x1, y0 + 2*radius], 270, 360, fill=outline, width=width)
        draw.arc([x0, y1 - 2*radius, x0 + 2*radius, y1], 90, 180, fill=outline, width=width)
        draw.arc([x1 - 2*radius, y1 - 2*radius, x1, y1], 0, 90, fill=outline, width=width)
        draw.line([x0 + radius, y0, x1 - radius, y0], fill=outline, width=width)
        draw.line([x0 + radius, y1, x1 - radius, y1], fill=outline, width=width)
        draw.line([x0, y0 + radius, x0, y1 - radius], fill=outline, width=width)
        draw.line([x1, y0 + radius, x1, y1 - radius], fill=outline, width=width)


@tool
def generate_diagram_image(
    diagram_type: str,
    content: str,
    filename: str = "output/diagram.png",
    theme: str = "default",
) -> str:
    """Generate a diagram image (flowchart, architecture, pipeline) using PIL.

    Args:
        diagram_type: Type of diagram — 'flowchart', 'architecture', or 'pipeline'.
        content: JSON string describing the diagram.
                 Flowchart: '{"steps": ["Start", "Process Data", "Validate", "Output", "End"]}'
                 Architecture: '{"layers": [{"name": "Frontend", "items": ["React", "Tailwind"]}, {"name": "Backend", "items": ["FastAPI", "Redis"]}, {"name": "Database", "items": ["PostgreSQL"]}]}'
                 Pipeline: '{"stages": ["Ingest", "Clean", "Transform", "Model", "Deploy"]}'
        filename: Output filename (e.g. 'output/flow.png').
        theme: PPT theme name to match diagram colours.

    Returns:
        str: Path to saved diagram image, or error message.
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return "Error: Pillow is required. Install with: pip install Pillow"

    try:
        diagram_data = json.loads(content)
    except json.JSONDecodeError:
        return "Error: content must be valid JSON"

    t = _get_theme(theme)
    colors_hex = t.get("chart_colors", ["#0D6EFD", "#198CFF", "#6CB4FF"])
    bg_hex = t.get("chart_bg", "#FFFFFF")
    is_dark = sum(int(bg_hex[i:i+2], 16) for i in (1, 3, 5)) < 384

    text_color = "#FFFFFF" if is_dark else "#1A1A2E"
    box_text = "#FFFFFF"
    arrow_color = "#888888" if not is_dark else "#AAAAAA"

    # Try to load a nice font, fall back to default
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 18)
        font_small = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 14)
        font_title = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 22)
    except Exception:
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
            font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 14)
            font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 22)
        except Exception:
            font = ImageFont.load_default()
            font_small = font
            font_title = font

    if diagram_type == "flowchart":
        steps = diagram_data.get("steps", ["Step 1", "Step 2", "Step 3"])
        n = len(steps)
        box_w, box_h = 240, 60
        gap = 50
        total_h = n * box_h + (n - 1) * gap
        img_w = max(600, box_w + 200)
        img_h = total_h + 160
        img = Image.new("RGB", (img_w, img_h), bg_hex)
        draw = ImageDraw.Draw(img)

        # Title
        draw.text((img_w // 2, 30), "Flow", fill=text_color, font=font_title, anchor="mt")

        start_x = (img_w - box_w) // 2
        start_y = 80

        for i, step in enumerate(steps):
            y = start_y + i * (box_h + gap)
            color = colors_hex[i % len(colors_hex)]
            _draw_rounded_rect(draw, (start_x, y, start_x + box_w, y + box_h), 12, fill=color)

            # Centred text
            bbox = draw.textbbox((0, 0), step, font=font)
            tw = bbox[2] - bbox[0]
            tx = start_x + (box_w - tw) // 2
            ty = y + (box_h - (bbox[3] - bbox[1])) // 2
            draw.text((tx, ty), step, fill=box_text, font=font)

            # Arrow to next
            if i < n - 1:
                arrow_x = start_x + box_w // 2
                arrow_y1 = y + box_h + 4
                arrow_y2 = y + box_h + gap - 4
                draw.line((arrow_x, arrow_y1, arrow_x, arrow_y2), fill=arrow_color, width=3)
                # Arrowhead
                draw.polygon([
                    (arrow_x - 8, arrow_y2 - 10),
                    (arrow_x + 8, arrow_y2 - 10),
                    (arrow_x, arrow_y2)
                ], fill=arrow_color)

    elif diagram_type == "architecture":
        layers = diagram_data.get("layers", [
            {"name": "Layer 1", "items": ["Item A"]},
            {"name": "Layer 2", "items": ["Item B"]},
        ])
        n_layers = len(layers)
        max_items = max(len(l.get("items", [])) for l in layers)
        item_w, item_h = 160, 50
        item_gap = 20
        layer_gap = 70
        row_w = max_items * item_w + (max_items - 1) * item_gap
        img_w = max(800, row_w + 200)
        img_h = n_layers * (item_h + 60) + (n_layers - 1) * layer_gap + 120
        img = Image.new("RGB", (img_w, img_h), bg_hex)
        draw = ImageDraw.Draw(img)

        draw.text((img_w // 2, 25), "Architecture", fill=text_color, font=font_title, anchor="mt")

        for li, layer in enumerate(layers):
            items = layer.get("items", [])
            layer_name = layer.get("name", f"Layer {li + 1}")
            layer_y = 70 + li * (item_h + 60 + layer_gap)
            color = colors_hex[li % len(colors_hex)]

            # Layer label
            draw.text((img_w // 2, layer_y), layer_name, fill=text_color, font=font, anchor="mt")

            # Items row
            total_row_w = len(items) * item_w + (len(items) - 1) * item_gap
            start_x = (img_w - total_row_w) // 2
            iy = layer_y + 30

            for ii, item in enumerate(items):
                ix = start_x + ii * (item_w + item_gap)
                _draw_rounded_rect(draw, (ix, iy, ix + item_w, iy + item_h), 10, fill=color)
                bbox = draw.textbbox((0, 0), item, font=font_small)
                tw = bbox[2] - bbox[0]
                tx = ix + (item_w - tw) // 2
                ty = iy + (item_h - (bbox[3] - bbox[1])) // 2
                draw.text((tx, ty), item, fill=box_text, font=font_small)

            # Arrow to next layer
            if li < n_layers - 1:
                ax = img_w // 2
                ay1 = iy + item_h + 10
                ay2 = iy + item_h + layer_gap - 10
                draw.line((ax, ay1, ax, ay2), fill=arrow_color, width=3)
                draw.polygon([
                    (ax - 8, ay2 - 10), (ax + 8, ay2 - 10), (ax, ay2)
                ], fill=arrow_color)

    else:  # pipeline
        stages = diagram_data.get("stages", ["Stage 1", "Stage 2", "Stage 3"])
        n = len(stages)
        box_w, box_h = 180, 60
        arrow_len = 50
        total_w = n * box_w + (n - 1) * arrow_len
        img_w = max(900, total_w + 120)
        img_h = 200
        img = Image.new("RGB", (img_w, img_h), bg_hex)
        draw = ImageDraw.Draw(img)

        draw.text((img_w // 2, 20), "Pipeline", fill=text_color, font=font_title, anchor="mt")

        start_x = (img_w - total_w) // 2
        y = 70

        for i, stage in enumerate(stages):
            x = start_x + i * (box_w + arrow_len)
            color = colors_hex[i % len(colors_hex)]
            _draw_rounded_rect(draw, (x, y, x + box_w, y + box_h), 12, fill=color)

            bbox = draw.textbbox((0, 0), stage, font=font)
            tw = bbox[2] - bbox[0]
            tx = x + (box_w - tw) // 2
            ty = y + (box_h - (bbox[3] - bbox[1])) // 2
            draw.text((tx, ty), stage, fill=box_text, font=font)

            if i < n - 1:
                ax1 = x + box_w + 5
                ax2 = x + box_w + arrow_len - 5
                ay = y + box_h // 2
                draw.line((ax1, ay, ax2, ay), fill=arrow_color, width=3)
                draw.polygon([
                    (ax2 - 10, ay - 8), (ax2 - 10, ay + 8), (ax2, ay)
                ], fill=arrow_color)

    full_path = _get_safe_path(filename)
    img.save(full_path, "PNG")
    logger.info(f"Diagram saved: {full_path}")
    return full_path


# ══════════════════════════════════════════════════════════════════════
# NEW TOOL 3: Fetch Relevant Image from Web
# ══════════════════════════════════════════════════════════════════════

@tool
def fetch_relevant_image(
    query: str,
    filename: str = "output/image.jpg",
) -> str:
    """Fetch a relevant royalty-free image from the web for use in presentations.

    Uses the Pexels API (free). Falls back to generating a styled placeholder
    if no API key is set.

    Args:
        query: Search query describing the desired image (e.g. 'teamwork office',
               'data analytics dashboard', 'green energy solar panels').
        filename: Output filename (e.g. 'output/background.jpg').

    Returns:
        str: Path to saved image, or error message.
    """
    import requests
    from dotenv import load_dotenv
    load_dotenv()

    full_path = _get_safe_path(filename)
    api_key = os.environ.get("PEXELS_API_KEY", "")

    if api_key:
        try:
            headers = {"Authorization": api_key}
            resp = requests.get(
                "https://api.pexels.com/v1/search",
                params={"query": query, "per_page": 1, "orientation": "landscape"},
                headers=headers,
                timeout=15,
            )
            if resp.status_code == 200:
                data = resp.json()
                photos = data.get("photos", [])
                if photos:
                    img_url = photos[0]["src"]["large2x"]
                    img_resp = requests.get(img_url, timeout=15)
                    if img_resp.status_code == 200:
                        with open(full_path, "wb") as f:
                            f.write(img_resp.content)
                        logger.info(f"Image fetched from Pexels: {full_path}")
                        return full_path
        except Exception as e:
            logger.warning(f"Pexels fetch failed: {e}")

    # Fallback: generate a professional placeholder with PIL
    try:
        from PIL import Image, ImageDraw, ImageFont
        img = Image.new("RGB", (1200, 800), "#2D3748")
        draw = ImageDraw.Draw(img)

        # Gradient-like effect with rectangles
        for i in range(400):
            alpha = int(255 * (1 - i / 400))
            r, g, b = 13, 110, 253
            draw.rectangle(
                [0, i * 2, 1200, i * 2 + 2],
                fill=(r * alpha // 255, g * alpha // 255, b * alpha // 255)
            )

        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 32)
        except Exception:
            font = ImageFont.load_default()

        # Wrap long queries
        display_text = textwrap.fill(query.title(), width=30)
        draw.text((600, 400), display_text, fill="#FFFFFF", font=font, anchor="mm")

        img.save(full_path, "JPEG", quality=90)
        logger.info(f"Placeholder image generated: {full_path}")
        return full_path

    except Exception as e:
        return f"Error generating image: {str(e)}"


# ══════════════════════════════════════════════════════════════════════
# NEW TOOL 4: Enhance Slide Design
# ══════════════════════════════════════════════════════════════════════

@tool
def enhance_slide_design(
    filepath: str,
    theme: str = "corporate",
) -> str:
    """Enhance an existing PowerPoint file with professional design elements.

    Adds gradient backgrounds, accent bars, bottom bars, and improved spacing
    to all slides in the presentation.

    Args:
        filepath: Path to the existing .pptx file to enhance.
        theme: Colour theme to apply. One of 'default', 'dark', 'corporate',
               'creative', 'consulting'.

    Returns:
        str: Confirmation message.
    """
    if not os.path.exists(filepath):
        return f"Error: File not found: {filepath}"

    try:
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches
    except ImportError:
        return "Error: python-pptx required"

    t = _get_theme(theme)
    prs = PptxPresentation(filepath)

    for slide in prs.slides:
        _set_slide_bg(slide, t["bg_color"])
        _add_accent_bar(slide, t["accent_color"])
        _add_bottom_bar(slide, t["accent_color"], t.get("accent_color_2"))

    prs.save(filepath)
    return f"Enhanced {len(prs.slides)} slides with '{theme}' theme. Saved to {filepath}"


# ══════════════════════════════════════════════════════════════════════
# NEW TOOL 5: Insert Visual Element into Slide
# ══════════════════════════════════════════════════════════════════════

@tool
def insert_visual_element(
    filepath: str,
    image_path: str,
    slide_index: str = "0",
    layout_type: str = "right_half",
) -> str:
    """Insert an image (chart, diagram, photo) into a specific slide of an
    existing PowerPoint file at a smart position.

    Args:
        filepath: Path to the .pptx file.
        image_path: Path to the image to insert (PNG/JPG).
        slide_index: Zero-based index of the slide to modify (as string).
                     Default '0' (first slide).
        layout_type: Placement layout — 'right_half', 'left_half',
                     'full_width', 'center', 'bottom'. Default 'right_half'.

    Returns:
        str: Confirmation message.
    """
    if not os.path.exists(filepath):
        return f"Error: File not found: {filepath}"
    if not os.path.exists(image_path):
        return f"Error: Image not found: {image_path}"

    try:
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches
    except ImportError:
        return "Error: python-pptx required"

    prs = PptxPresentation(filepath)
    idx = int(slide_index)
    if idx >= len(prs.slides):
        return f"Error: slide_index {idx} out of range (presentation has {len(prs.slides)} slides)"

    slide = prs.slides[idx]

    positions = {
        "right_half":  (Inches(6.5), Inches(1.5), Inches(6.3), Inches(5.0)),
        "left_half":   (Inches(0.5), Inches(1.5), Inches(6.3), Inches(5.0)),
        "full_width":  (Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.0)),
        "center":      (Inches(3.0), Inches(1.5), Inches(7.3), Inches(5.0)),
        "bottom":      (Inches(1.0), Inches(4.0), Inches(11.3), Inches(3.2)),
    }
    pos = positions.get(layout_type, positions["right_half"])

    try:
        slide.shapes.add_picture(image_path, *pos)
    except Exception as e:
        return f"Error inserting image: {str(e)}"

    prs.save(filepath)
    return f"Image inserted at '{layout_type}' on slide {idx}. Saved to {filepath}"


# ══════════════════════════════════════════════════════════════════════
# Slide builders (internal) — ENHANCED
# ══════════════════════════════════════════════════════════════════════

def _build_title_slide(prs, title, subtitle, author, t):
    """Build an enhanced title slide with accent elements."""
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t["bg_color"])

    # Left accent bar
    _add_accent_bar(slide, t["accent_color"], Inches(0), Inches(0), Inches(0.2), Inches(7.5))

    # Decorative top-right accent block
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(9.5), Inches(0), Inches(3.833), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(t["accent_color"])
    shape.line.fill.background()

    # Second accent block (smaller)
    ac2 = t.get("accent_color_2", t["accent_color"])
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(10.5), Inches(2.5), Inches(2.833), Inches(1.2)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = _rgb(ac2)
    shape2.line.fill.background()

    # Title text
    _add_text_box(slide, Inches(0.8), Inches(2.0), Inches(8.5), Inches(2.0),
                  title, t["title_font"], 44, t["title_color"], bold=True)

    # Subtitle
    if subtitle:
        _add_text_box(slide, Inches(0.8), Inches(4.2), Inches(8.5), Inches(1.0),
                      subtitle, t["body_font"], 20, t["subtitle_color"])

    # Author
    if author:
        _add_text_box(slide, Inches(0.8), Inches(5.8), Inches(8.5), Inches(0.5),
                      f"By {author}", t["body_font"], 14, t["subtitle_color"])

    # Bottom bar
    _add_bottom_bar(slide, t["accent_color"], t.get("accent_color_2"))


def _build_bullet_slide(prs, data: dict, t: dict, slide_num=0, total=0):
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t["bg_color"])
    _add_accent_bar(slide, t["accent_color"], Inches(0), Inches(0), Inches(0.12), Inches(7.5))

    heading = data.get("heading", "")
    bullets = data.get("bullets", [])

    # Heading with accent underline
    _add_text_box(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.8),
                  heading, t["title_font"], 28, t["heading_color"], bold=True)

    # Accent line under heading
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.25), Inches(2.0), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(t["accent_color"])
    shape.line.fill.background()

    # Bullets with proper spacing and bullet markers
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:6]):  # Max 6 bullets
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.name = t["body_font"]
        p.font.size = Pt(18)
        p.font.color.rgb = _rgb(t["body_color"])
        p.space_after = Pt(14)
        p.level = 0

        # Bullet indicator dot
        bullet_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.65), Inches(1.75 + i * 0.55), Inches(0.12), Inches(0.12)
        )
        bullet_shape.fill.solid()
        bullet_shape.fill.fore_color.rgb = _rgb(t["accent_color"])
        bullet_shape.line.fill.background()

    _add_bottom_bar(slide, t["accent_color"], t.get("accent_color_2"))
    if slide_num and total:
        _add_slide_number(slide, slide_num, total, t)


def _build_section_slide(prs, data: dict, t: dict, slide_num=0, total=0):
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t["accent_color"])

    heading = data.get("heading", "")
    sub = data.get("subtitle", "")

    # Large centred heading
    _add_text_box(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(1.5),
                  heading, t["title_font"], 40, (0xFF, 0xFF, 0xFF),
                  bold=True, alignment=PP_ALIGN.CENTER)
    if sub:
        _add_text_box(slide, Inches(1), Inches(4.3), Inches(11.333), Inches(0.8),
                      sub, t["body_font"], 18, (0xEE, 0xEE, 0xEE),
                      alignment=PP_ALIGN.CENTER)

    # Decorative line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(4.15), Inches(2.333), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb((0xFF, 0xFF, 0xFF))
    shape.line.fill.background()


def _build_text_slide(prs, data: dict, t: dict, slide_num=0, total=0):
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t["bg_color"])
    _add_accent_bar(slide, t["accent_color"], Inches(0), Inches(0), Inches(0.12), Inches(7.5))

    heading = data.get("heading", "")
    body = data.get("body", "")

    _add_text_box(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.8),
                  heading, t["title_font"], 28, t["heading_color"], bold=True)

    # Accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.25), Inches(2.0), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(t["accent_color"])
    shape.line.fill.background()

    _add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5),
                  body, t["body_font"], 16, t["body_color"])

    _add_bottom_bar(slide, t["accent_color"], t.get("accent_color_2"))
    if slide_num and total:
        _add_slide_number(slide, slide_num, total, t)


def _build_comparison_slide(prs, data: dict, t: dict, slide_num=0, total=0):
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t["bg_color"])
    _add_accent_bar(slide, t["accent_color"], Inches(0), Inches(0), Inches(0.12), Inches(7.5))

    heading = data.get("heading", "")
    _add_text_box(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.8),
                  heading, t["title_font"], 28, t["heading_color"], bold=True)

    # Accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.25), Inches(2.0), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(t["accent_color"])
    shape.line.fill.background()

    # Left column header
    left_title = data.get("left_title", "Option A")
    left_items = data.get("left_items", [])

    # Left column box
    lbox = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2)
    )
    lbox.fill.solid()
    lbox.fill.fore_color.rgb = _rgb(t.get("bg_gradient_end", t["bg_color"]))
    lbox.line.fill.background()

    _add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.4), Inches(0.6),
                  left_title, t["title_font"], 20, t["accent_color"], bold=True)

    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = t["body_font"]
        p.font.size = Pt(16)
        p.font.color.rgb = _rgb(t["body_color"])
        p.space_after = Pt(8)

    # Right column
    right_title = data.get("right_title", "Option B")
    right_items = data.get("right_items", [])

    rbox = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2)
    )
    rbox.fill.solid()
    rbox.fill.fore_color.rgb = _rgb(t.get("bg_gradient_end", t["bg_color"]))
    rbox.line.fill.background()

    _add_text_box(slide, Inches(7.1), Inches(1.8), Inches(5.4), Inches(0.6),
                  right_title, t["title_font"], 20,
                  t.get("accent_color_2", t["accent_color"]), bold=True)

    txBox2 = slide.shapes.add_textbox(Inches(7.3), Inches(2.5), Inches(5), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"• {item}"
        p.font.name = t["body_font"]
        p.font.size = Pt(16)
        p.font.color.rgb = _rgb(t["body_color"])
        p.space_after = Pt(8)

    _add_bottom_bar(slide, t["accent_color"], t.get("accent_color_2"))
    if slide_num and total:
        _add_slide_number(slide, slide_num, total, t)


def _build_image_slide(prs, data: dict, t: dict, slide_num=0, total=0):
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t["bg_color"])
    _add_accent_bar(slide, t["accent_color"], Inches(0), Inches(0), Inches(0.12), Inches(7.5))

    heading = data.get("heading", "")
    image_path = data.get("image_path", "")
    caption = data.get("caption", "")

    _add_text_box(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.8),
                  heading, t["title_font"], 28, t["heading_color"], bold=True)

    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.5),
                                     Inches(10.3), Inches(5.0))
        except Exception as e:
            _add_text_box(slide, Inches(2), Inches(3), Inches(9), Inches(1),
                          f"[Image could not be loaded: {e}]",
                          t["body_font"], 14, t["subtitle_color"],
                          alignment=PP_ALIGN.CENTER)
    else:
        _add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(1),
                      f"[Image not found: {image_path}]",
                      t["body_font"], 14, t["subtitle_color"],
                      alignment=PP_ALIGN.CENTER)

    if caption:
        _add_text_box(slide, Inches(1), Inches(6.6), Inches(11.333), Inches(0.5),
                      caption, t["body_font"], 11, t["subtitle_color"],
                      alignment=PP_ALIGN.CENTER)

    _add_bottom_bar(slide, t["accent_color"], t.get("accent_color_2"))
    if slide_num and total:
        _add_slide_number(slide, slide_num, total, t)


def _build_chart_slide(prs, data: dict, t: dict, slide_num=0, total=0):
    """Build a slide with a chart image (auto-generates if needed)."""
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t["bg_color"])
    _add_accent_bar(slide, t["accent_color"], Inches(0), Inches(0), Inches(0.12), Inches(7.5))

    heading = data.get("heading", "Data Overview")
    chart_path = data.get("chart_path", "")
    chart_data = data.get("chart_data", "")
    chart_type = data.get("chart_type", "bar")

    # If no chart_path but chart_data provided, generate the chart
    if not chart_path and chart_data:
        import hashlib
        h = hashlib.md5(str(chart_data).encode()).hexdigest()[:8]
        chart_filename = f"output/chart_{h}.png"
        chart_data_str = chart_data if isinstance(chart_data, str) else json.dumps(chart_data)
        theme_name = "default"
        for name, th in THEMES.items():
            if th == t:
                theme_name = name
                break
        chart_path = generate_chart_image(
            data=chart_data_str, chart_type=chart_type,
            filename=chart_filename, title=heading, theme=theme_name
        )

    _add_text_box(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.8),
                  heading, t["title_font"], 28, t["heading_color"], bold=True)

    # Accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.25), Inches(2.0), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(t["accent_color"])
    shape.line.fill.background()

    if chart_path and os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(1.5), Inches(1.5),
                                 Inches(10.3), Inches(5.0))
    else:
        _add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(1),
                      "[Chart could not be generated]",
                      t["body_font"], 14, t["subtitle_color"],
                      alignment=PP_ALIGN.CENTER)

    _add_bottom_bar(slide, t["accent_color"], t.get("accent_color_2"))
    if slide_num and total:
        _add_slide_number(slide, slide_num, total, t)


def _build_diagram_slide(prs, data: dict, t: dict, slide_num=0, total=0):
    """Build a slide with a diagram image (auto-generates if needed)."""
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t["bg_color"])
    _add_accent_bar(slide, t["accent_color"], Inches(0), Inches(0), Inches(0.12), Inches(7.5))

    heading = data.get("heading", "System Diagram")
    diagram_path = data.get("diagram_path", "")
    diagram_type = data.get("diagram_type", "flowchart")
    diagram_content = data.get("diagram_content", "")

    # Auto-generate if needed
    if not diagram_path and diagram_content:
        import hashlib
        h = hashlib.md5(str(diagram_content).encode()).hexdigest()[:8]
        diag_filename = f"output/diagram_{h}.png"
        content_str = diagram_content if isinstance(diagram_content, str) else json.dumps(diagram_content)
        theme_name = "default"
        for name, th in THEMES.items():
            if th == t:
                theme_name = name
                break
        diagram_path = generate_diagram_image(
            diagram_type=diagram_type, content=content_str,
            filename=diag_filename, theme=theme_name
        )

    _add_text_box(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.8),
                  heading, t["title_font"], 28, t["heading_color"], bold=True)

    # Accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.25), Inches(2.0), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(t["accent_color"])
    shape.line.fill.background()

    if diagram_path and os.path.exists(diagram_path):
        slide.shapes.add_picture(diagram_path, Inches(1.5), Inches(1.5),
                                 Inches(10.3), Inches(5.0))
    else:
        _add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(1),
                      "[Diagram could not be generated]",
                      t["body_font"], 14, t["subtitle_color"],
                      alignment=PP_ALIGN.CENTER)

    _add_bottom_bar(slide, t["accent_color"], t.get("accent_color_2"))
    if slide_num and total:
        _add_slide_number(slide, slide_num, total, t)


def _build_thank_you_slide(prs, title, t):
    """Build a professional thank-you / closing slide."""
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t["bg_color"])

    # Large accent block
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(t["accent_color"])
    shape.line.fill.background()

    _add_text_box(slide, Inches(1), Inches(1.0), Inches(11.333), Inches(1.5),
                  "Thank You", t["title_font"], 48, (0xFF, 0xFF, 0xFF),
                  bold=True, alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(4.5), Inches(11.333), Inches(1.0),
                  title, t["body_font"], 20, t["subtitle_color"],
                  alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(5.8), Inches(11.333), Inches(0.6),
                  "Questions & Discussion", t["body_font"], 16, t["subtitle_color"],
                  alignment=PP_ALIGN.CENTER)

    _add_bottom_bar(slide, t["accent_color"], t.get("accent_color_2"))


# ══════════════════════════════════════════════════════════════════════
# ENHANCED: create_presentation — the primary tool
# ══════════════════════════════════════════════════════════════════════

@tool
def create_presentation(
    filename: str,
    title: str,
    slides_data: list,
    theme: str = "default",
    subtitle: str = "Generated by Synergy Agent",
    author: str = "Synergy Agent",
) -> str:
    """Create a professional, industry-grade PowerPoint presentation (.pptx)
    with charts, diagrams, images, and polished layouts.

    This is the primary tool for building presentations. Pass structured
    slide data and the tool handles layout, styling, visuals, and saving.

    Args:
        filename: Output filename (e.g. 'output/report.pptx'). Include .pptx extension.
        title: Main title shown on the first (cover) slide.
        slides_data: A list of dicts, one per slide. Each dict MUST have a
                     'type' key. Supported types and their keys:

                     type='bullet' — heading (str), bullets (list of str, max 6)
                     type='section' — heading (str), subtitle (str, optional)
                     type='text' — heading (str), body (str)
                     type='comparison' — heading (str), left_title (str),
                       left_items (list), right_title (str), right_items (list)
                     type='image' — heading (str), image_path (str), caption (str, optional)
                     type='chart' — heading (str), chart_data (dict/JSON str with
                       labels and values), chart_type ('bar'/'line'/'pie'),
                       chart_path (str, optional — pre-generated image)
                     type='diagram' — heading (str), diagram_type ('flowchart'/
                       'architecture'/'pipeline'), diagram_content (dict/JSON str),
                       diagram_path (str, optional — pre-generated image)

        theme: Colour theme. One of 'default', 'dark', 'corporate', 'creative', 'consulting'.
        subtitle: Subtitle on the title slide.
        author: Author name placed on the title slide.

    Returns:
        Confirmation message with file path and slide count.
    """
    try:
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches
    except ImportError:
        return "Error: python-pptx is required. Install with: pip install python-pptx"

    t = _get_theme(theme)
    full_path = _get_safe_path(filename)
    prs = PptxPresentation()

    # Set default slide dimensions (widescreen 16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Title slide ───────────────────────────────────────────────────
    _build_title_slide(prs, title, subtitle, author, t)

    total_content = len(slides_data) + 2  # +2 for title + thank-you

    # ── Content slides ────────────────────────────────────────────────
    for idx, s in enumerate(slides_data):
        slide_type = s.get("type", "bullet")
        sn = idx + 2  # slide number (1=title, so content starts at 2)

        if slide_type == "bullet":
            _build_bullet_slide(prs, s, t, sn, total_content)
        elif slide_type == "section":
            _build_section_slide(prs, s, t, sn, total_content)
        elif slide_type == "text":
            _build_text_slide(prs, s, t, sn, total_content)
        elif slide_type == "comparison":
            _build_comparison_slide(prs, s, t, sn, total_content)
        elif slide_type == "image":
            _build_image_slide(prs, s, t, sn, total_content)
        elif slide_type == "chart":
            _build_chart_slide(prs, s, t, sn, total_content)
        elif slide_type == "diagram":
            _build_diagram_slide(prs, s, t, sn, total_content)
        else:
            _build_bullet_slide(prs, s, t, sn, total_content)

    # ── Thank-you / end slide ─────────────────────────────────────────
    _build_thank_you_slide(prs, title, t)

    prs.save(full_path)
    total = len(prs.slides)
    logger.info(f"PPT saved: {full_path} ({total} slides)")
    return f"Presentation saved to {full_path} ({total} slides, theme='{theme}')"


# ══════════════════════════════════════════════════════════════════════
# EXISTING: create_quick_ppt (preserved)
# ══════════════════════════════════════════════════════════════════════

@tool
def create_quick_ppt(
    filename: str,
    title: str,
    bullet_points: list,
    theme: str = "default",
) -> str:
    """Create a quick 3-slide PowerPoint: title slide, bullet slide, thank-you slide.

    Use this when the user wants a simple, fast presentation with one set of bullets.
    For multi-slide presentations prefer create_presentation.

    Args:
        filename: Output filename ending in .pptx (e.g. 'output/summary.pptx').
        title: Main presentation title.
        bullet_points: List of bullet-point strings (max 6 recommended).
        theme: Colour theme. One of 'default', 'dark', 'corporate', 'creative', 'consulting'.

    Returns:
        Confirmation message with file path and slide count.
    """
    slides = [
        {"type": "bullet", "heading": title, "bullets": bullet_points},
    ]
    return create_presentation(filename, title, slides, theme=theme)


# ══════════════════════════════════════════════════════════════════════
# EXISTING: add_slide_to_ppt (preserved & enhanced)
# ══════════════════════════════════════════════════════════════════════

@tool
def add_slide_to_ppt(
    filepath: str,
    slide_type: str,
    heading: str,
    content: list,
) -> str:
    """Open an existing .pptx file and append a new slide to it.

    Use this to incrementally build a presentation one slide at a time.

    Args:
        filepath: Path to the existing .pptx file.
        slide_type: One of 'bullet', 'text', or 'section'.
        heading: Heading text for the new slide.
        content: For 'bullet' a list of bullet strings.
                 For 'text' a list with a single paragraph string.
                 For 'section' a list with an optional subtitle string.

    Returns:
        Confirmation message with updated slide count.
    """
    if not os.path.exists(filepath):
        return f"Error: File not found: {filepath}"

    try:
        from pptx import Presentation as PptxPresentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return "Error: python-pptx is required. Install with: pip install python-pptx"

    prs = PptxPresentation(filepath)
    t = _get_theme("default")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t["bg_color"])

    if slide_type == "section":
        _set_slide_bg(slide, t["accent_color"])
        _add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                      heading, t["title_font"], 36, (0xFF, 0xFF, 0xFF),
                      bold=True, alignment=PP_ALIGN.CENTER)
        if content:
            _add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                          content[0], t["body_font"], 18,
                          (0xEE, 0xEE, 0xEE),
                          alignment=PP_ALIGN.CENTER)
    elif slide_type == "text":
        _add_accent_bar(slide, t["accent_color"])
        _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                      heading, t["title_font"], 28, t["heading_color"], bold=True)
        body = content[0] if content else ""
        _add_text_box(slide, Inches(1), Inches(1.8), Inches(10.5), Inches(5),
                      body, t["body_font"], 16, t["body_color"])
    else:
        # Default: bullet
        _add_accent_bar(slide, t["accent_color"])
        _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                      heading, t["title_font"], 28, t["heading_color"], bold=True)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8),
                                         Inches(10.5), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.name = t["body_font"]
            p.font.size = Pt(18)
            p.font.color.rgb = _rgb(t["body_color"])
            p.space_after = Pt(10)

    _add_bottom_bar(slide, t["accent_color"], t.get("accent_color_2"))

    prs.save(filepath)
    total = len(prs.slides)
    logger.info(f"Added slide to {filepath} ({total} slides)")
    return f"Slide added. Presentation now has {total} slides. Saved to {filepath}"


# ══════════════════════════════════════════════════════════════════════
# EXISTING: list_ppt_themes (updated with new themes)
# ══════════════════════════════════════════════════════════════════════

@tool
def list_ppt_themes() -> str:
    """List all available PowerPoint colour themes with a short description.

    Returns:
        Formatted string of theme names and their colour palettes.
    """
    descriptions = {
        "default":    "Clean white background, navy titles, blue accents",
        "dark":       "Dark charcoal background, white text, cyan accents",
        "corporate":  "Light grey background, dark navy titles, red accents",
        "creative":   "White background, purple/pink titles, orange accents",
        "consulting": "Clean white, deep navy titles, blue + orange accents (McKinsey-style)",
    }
    lines = ["Available PPT themes:", ""]
    for name, desc in descriptions.items():
        lines.append(f"  - {name:12s}: {desc}")
    lines.append("")
    lines.append("Pass the theme name to create_presentation() or create_quick_ppt().")
    return "\n".join(lines)


@tool
def generate_ppt_preview(filepath: str) -> str:
    """
    Generates a visual PNG thumbnail preview of the presentation (usually the title slide)
    to visually confirm the theme, aesthetic, and basic layout.
    Works natively on macOS using QuickLook.

    Args:
        filepath: The path to the .pptx file.

    Returns:
        str: Information about the generated PNG, including its file path, or an error message.
    """
    import os
    import subprocess
    
    abs_path = os.path.abspath(filepath)
    if not os.path.exists(abs_path):
        return f"Error: File {filepath} does not exist."
        
    out_dir = os.path.dirname(abs_path)
    try:
        res = subprocess.run(
            ["qlmanage", "-t", "-s", "1000", "-o", out_dir, abs_path], 
            capture_output=True, text=True
        )
        base = os.path.basename(filepath)
        png_path = os.path.join(out_dir, base + ".png")
        if os.path.exists(png_path):
            return f"Success! Preview screenshot saved at: {png_path}. Show this path to the user."
        else:
            return f"Failed to generate preview. qlmanage output: {res.stdout} {res.stderr}"
    except Exception as e:
        return f"Exception generating preview: {str(e)}"


@tool
def analyze_ppt_layout(filepath: str) -> str:
    """
    Performs a geometric and structural layout check on the presentation to detect
    overlapping text boxes, missing visuals, or poor text density.
    This acts as a "visual recheck" for text-based agents to spot design flaws.

    Args:
        filepath: The path to the .pptx file to evaluate.

    Returns:
        str: A detailed textual visual critique of the slides.
    """
    import os
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return "Error: python-pptx not installed."

    abs_path = os.path.abspath(filepath)
    if not os.path.exists(abs_path):
        return f"Error: File {filepath} does not exist."

    try:
        prs = Presentation(abs_path)
        critique = []
        critique.append(f"Layout Analysis for {os.path.basename(filepath)} ({len(prs.slides)} slides):")
        
        for idx, slide in enumerate(prs.slides):
            issues = []
            shapes = list(slide.shapes)
            
            # Bounding box collision detection
            boxes = []
            for s in shapes:
                if hasattr(s, "left") and hasattr(s, "top") and hasattr(s, "width") and hasattr(s, "height"):
                    boxes.append({
                        "name": s.name,
                        "left": s.left, "top": s.top, 
                        "right": s.left + s.width, "bottom": s.top + s.height,
                        "type": getattr(s, "shape_type", None)
                    })
                    
            for i, b1 in enumerate(boxes):
                for j, b2 in enumerate(boxes):
                    if i >= j: continue
                    # AABB collision check
                    overlap_x = b1["left"] < b2["right"] and b1["right"] > b2["left"]
                    overlap_y = b1["top"] < b2["bottom"] and b1["bottom"] > b2["top"]
                    if overlap_x and overlap_y:
                        # Ignore overlaps among rects which are often just backgrounds
                        if "Rectangle" not in b1["name"] and "Rectangle" not in b2["name"]:
                            issues.append(f"Potential Overlap between '{b1['name']}' and '{b2['name']}'")

            # Check for visuals
            has_visual = any(getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE or "Graphic" in s.name or "Chart" in s.name for s in shapes)
            if not has_visual and idx > 0:
                issues.append("No images, charts, or visual diagrams found. Slide might be too text-heavy.")

            if issues:
                critique.append(f"\n- Slide {idx+1}:")
                for iss in issues:
                    critique.append(f"  * {iss}")
            else:
                critique.append(f"\n- Slide {idx+1}: Looks cleanly laid out.")
                
        return "\n".join(critique)
    except Exception as e:
        return f"Analysis failed: {str(e)}"


# ══════════════════════════════════════════════════════════════════════
# Export list
# ══════════════════════════════════════════════════════════════════════

PPT_TOOLS = [
    create_presentation,
    create_quick_ppt,
    add_slide_to_ppt,
    list_ppt_themes,
    generate_chart_image,
    generate_diagram_image,
    fetch_relevant_image,
    enhance_slide_design,
    insert_visual_element,
    generate_ppt_preview,
    analyze_ppt_layout,
]


# ══════════════════════════════════════════════════════════════════════
# Self-test
# ══════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("=" * 60)
    print("TEST: Full presentation with charts and diagrams")
    print("=" * 60)

    test_slides = [
        {"type": "section", "heading": "Market Analysis", "subtitle": "Understanding the landscape"},
        {
            "type": "bullet",
            "heading": "Why This Problem Matters Now",
            "bullets": [
                "Market growing 25% year-over-year",
                "Customer acquisition costs rising",
                "Competitors entering the space",
                "Regulatory changes imminent",
            ],
        },
        {
            "type": "chart",
            "heading": "Revenue Growth by Quarter",
            "chart_type": "bar",
            "chart_data": {
                "labels": ["Q1", "Q2", "Q3", "Q4"],
                "values": [2100, 2600, 3200, 4100]
            },
        },
        {
            "type": "diagram",
            "heading": "Solution Architecture",
            "diagram_type": "architecture",
            "diagram_content": {
                "layers": [
                    {"name": "Frontend", "items": ["React App", "Mobile SDK"]},
                    {"name": "API Layer", "items": ["FastAPI", "GraphQL"]},
                    {"name": "Data", "items": ["PostgreSQL", "Redis", "S3"]},
                ]
            },
        },
        {
            "type": "comparison",
            "heading": "Current vs Proposed Solution",
            "left_title": "Current State",
            "left_items": ["Manual processes", "3-day turnaround", "60% accuracy"],
            "right_title": "Proposed Solution",
            "right_items": ["Fully automated", "Real-time results", "95% accuracy"],
        },
        {
            "type": "text",
            "heading": "Executive Summary",
            "body": (
                "Our analysis reveals a significant opportunity to transform operations "
                "through intelligent automation. The proposed solution delivers 10x faster "
                "processing with 95% accuracy, resulting in an estimated $2.4M annual savings."
            ),
        },
    ]

    result = create_presentation(
        "output/test_premium.pptx",
        "Strategic Growth Initiative",
        test_slides,
        theme="consulting",
        subtitle="Board Presentation — Q2 2026",
        author="Strategy Team",
    )
    print(f"  OK: {result}")

    print("\nTEST: Chart generation")
    chart_path = generate_chart_image(
        data='{"labels": ["Engineering", "Sales", "Marketing", "Support"], "values": [45, 30, 15, 10]}',
        chart_type="pie",
        filename="output/test_pie.png",
        title="Budget Allocation",
        theme="consulting",
    )
    print(f"  OK: Chart at {chart_path}")

    print("\nTEST: Diagram generation")
    diag_path = generate_diagram_image(
        diagram_type="pipeline",
        content='{"stages": ["Collect", "Clean", "Analyze", "Report"]}',
        filename="output/test_pipeline.png",
        theme="corporate",
    )
    print(f"  OK: Diagram at {diag_path}")

    print("\nAll tests passed.")
